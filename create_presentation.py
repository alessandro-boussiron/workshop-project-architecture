#!/usr/bin/env python3
"""
Générateur de présentation PowerPoint pour le Workshop Project Architecture
Version courte et visuelle
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_title_slide(prs, title, subtitle=""):
    """Crée une diapositive de titre"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    
    if subtitle:
        subtitle_shape.text = subtitle
        subtitle_shape.text_frame.paragraphs[0].font.size = Pt(28)
    
    return slide

def create_content_slide(prs, title, content_items, font_size=20):
    """Crée une diapositive avec titre et contenu bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    title_shape.text = title
    
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    
    for item in content_items:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(font_size)
        p.space_after = Pt(12)
    
    return slide

def create_two_column_slide(prs, title, left_items, right_items):
    """Crée une diapositive avec deux colonnes"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    # Titre
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Colonne gauche
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
    left_frame = left_box.text_frame
    for item in left_items:
        p = left_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
    
    # Colonne droite
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.5), Inches(5))
    right_frame = right_box.text_frame
    for item in right_items:
        p = right_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
    
    return slide

def add_diagram_box(slide, left, top, width, height, text, fill_color=None, font_size=16):
    """Ajoute une boîte avec texte pour les diagrammes"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    
    shape.shadow.inherit = False
    
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    
    # Bordure
    shape.line.color.rgb = RGBColor(50, 50, 50)
    shape.line.width = Pt(2)
    
    text_frame = shape.text_frame
    text_frame.text = text
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_frame.paragraphs[0].font.size = Pt(font_size)
    text_frame.paragraphs[0].font.bold = True
    text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    text_frame.word_wrap = True
    text_frame.vertical_anchor = 1  # Middle
    
    return shape

def add_arrow(slide, x1, y1, x2, y2):
    """Ajoute une flèche entre deux points"""
    connector = slide.shapes.add_connector(
        2,  # Straight connector
        x1, y1, x2, y2
    )
    connector.line.color.rgb = RGBColor(50, 50, 50)
    connector.line.width = Pt(3)
    return connector

def add_big_icon_text(slide, top, emoji, text, color):
    """Ajoute un grand emoji avec texte"""
    # Emoji
    emoji_box = slide.shapes.add_textbox(Inches(3), top, Inches(4), Inches(1))
    emoji_frame = emoji_box.text_frame
    emoji_frame.text = emoji
    p = emoji_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(72)
    
    # Texte
    text_box = slide.shapes.add_textbox(Inches(1.5), top + Inches(1), Inches(7), Inches(1))
    text_frame = text_box.text_frame
    text_frame.text = text
    p = text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = color
    
    return emoji_box, text_box

def create_architecture_diagram_slide(prs):
    """Crée une diapositive avec diagramme d'architecture en couches"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    title_shape = slide.shapes.title
    title_shape.text = "🏗️ Architecture en Couches"
    
    # Couches avec couleurs vives
    layers = [
        ("PRESENTATION\nHTTP • GraphQL • CLI", RGBColor(52, 152, 219)),
        ("APPLICATION\nUse Cases", RGBColor(46, 204, 113)),
        ("DOMAIN\nBusiness Logic", RGBColor(241, 196, 15)),
        ("INFRASTRUCTURE\nDB • Email • Files", RGBColor(231, 76, 60))
    ]
    
    top = Inches(1.8)
    for text, color in layers:
        add_diagram_box(slide, Inches(2.5), top, Inches(5), Inches(1.1), text, color, 18)
        top += Inches(1.3)
        
        # Flèche vers le bas (sauf pour le dernier)
        if text != layers[-1][0]:
            arrow_y = top - Inches(0.2)
            add_arrow(slide, Inches(5), arrow_y, Inches(5), arrow_y + Inches(0.2))
    
    # Note en bas
    note = slide.shapes.add_textbox(Inches(1), Inches(6.8), Inches(8), Inches(0.5))
    note.text_frame.text = "⬇️ Dépendances pointent vers l'intérieur"
    note.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    note.text_frame.paragraphs[0].font.size = Pt(20)
    note.text_frame.paragraphs[0].font.italic = True
    note.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)
    
    return slide

def create_cqrs_diagram_slide(prs):
    """Crée une diapositive avec diagramme CQRS"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    title_shape = slide.shapes.title
    title_shape.text = "⚡ CQRS Pattern"
    
    # Write Side
    add_diagram_box(slide, Inches(0.5), Inches(2), Inches(4), Inches(2), 
                    "📝 WRITE\nCommands\n\nCreate\nUpdate\nDelete", 
                    RGBColor(231, 76, 60), 20)
    
    # Read Side
    add_diagram_box(slide, Inches(5.5), Inches(2), Inches(4), Inches(2), 
                    "📖 READ\nQueries\n\nGet\nSearch\nFilter", 
                    RGBColor(52, 152, 219), 20)
    
    # Event Store au centre
    add_diagram_box(slide, Inches(2.5), Inches(5), Inches(5), Inches(1.2), 
                    "💾 Event Store / Database", 
                    RGBColor(149, 165, 166), 18)
    
    # Flèches
    add_arrow(slide, Inches(2.5), Inches(3), Inches(4.5), Inches(5))
    add_arrow(slide, Inches(7.5), Inches(4), Inches(5.5), Inches(5.5))
    
    return slide

def create_nasa_rules_slide(prs):
    """Crée une diapositive sur les règles NASA"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title_shape = slide.shapes.title
    title_shape.text = "NASA Power of 10 Rules"
    
    rules = [
        "1. Simple control flow (pas de goto, récursion)",
        "2. Fixed loop bounds (boucles bornées)",
        "3. No dynamic allocation (pas malloc après init)",
        "4. Small functions (< 60 lignes)",
        "5. Check return values (toujours vérifier)",
        "6. Limited scope (variables au plus petit scope)",
        "7. Check assertions (assert() partout)",
        "8. Limited pointers (max 2 niveaux)",
        "9. Zero warnings (-Wall -Wextra -Werror)",
        "10. Static analysis (Coverity, PC-Lint)"
    ]
    
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    
    for rule in rules:
        p = tf.add_paragraph()
        p.text = rule
        p.font.size = Pt(16)
    
    return slide

def create_memory_hierarchy_slide(prs):
    """Crée une diapositive sur la hiérarchie mémoire avec visuels"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    title_shape = slide.shapes.title
    title_shape.text = "🛡️ Hiérarchie Mémoire (SANS malloc!)"
    
    options = [
        ("⭐⭐⭐⭐⭐\nSTATIC\nZéro fuite", RGBColor(46, 204, 113), 1.8),
        ("⭐⭐⭐⭐\nOBJECT POOL\nRéutilisation", RGBColor(52, 152, 219), 3.2),
        ("⭐⭐⭐\nARENA\nSimple", RGBColor(241, 196, 15), 4.6),
        ("⭐⭐\nMALLOC\nDernier recours", RGBColor(231, 76, 60), 6.0)
    ]
    
    for text, color, top in options:
        add_diagram_box(slide, Inches(2), Inches(top), Inches(6), Inches(1.2), text, color, 18)
    
    return slide

def create_solid_visual_slide(prs):
    """Crée une diapositive SOLID visuelle"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    title_shape = slide.shapes.title
    title_shape.text = "🎯 Principes SOLID"
    
    principles = [
        ("S", "Single\nResponsibility", RGBColor(52, 152, 219), 0.8, 2),
        ("O", "Open/\nClosed", RGBColor(46, 204, 113), 2.8, 2),
        ("L", "Liskov\nSubstitution", RGBColor(241, 196, 15), 4.8, 2),
        ("I", "Interface\nSegregation", RGBColor(230, 126, 34), 6.8, 2),
        ("D", "Dependency\nInversion", RGBColor(155, 89, 182), 0.8, 4.5),
    ]
    
    for letter, text, color, left, top in principles[:4]:
        # Lettre
        box = add_diagram_box(slide, Inches(left), Inches(top), Inches(1.8), Inches(2), 
                             f"{letter}\n\n{text}", color, 16)
    
    # D plus large en bas
    add_diagram_box(slide, Inches(2.5), Inches(4.5), Inches(5), Inches(2), 
                   f"D\n\n{principles[4][1]}", principles[4][2], 18)
    
    return slide

def create_ddd_visual_slide(prs):
    """Crée une diapositive DDD très visuelle"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    title_shape = slide.shapes.title
    title_shape.text = "📦 DDD Building Blocks"
    
    blocks = [
        ("🎯\nEntity", RGBColor(52, 152, 219), 0.8, 2.5),
        ("💎\nValue\nObject", RGBColor(46, 204, 113), 3, 2.5),
        ("📦\nAggregate", RGBColor(241, 196, 15), 5.2, 2.5),
        ("⚡\nDomain\nEvent", RGBColor(230, 126, 34), 7.4, 2.5),
        ("🔧\nService", RGBColor(231, 76, 60), 1.9, 5),
        ("🏭\nFactory", RGBColor(155, 89, 182), 4.1, 5),
        ("📚\nRepository", RGBColor(52, 73, 94), 6.3, 5),
    ]
    
    for text, color, left, top in blocks:
        add_diagram_box(slide, Inches(left), Inches(top), Inches(2), Inches(1.5), 
                       text, color, 16)
    
    return slide

def create_nasa_visual_slide(prs):
    """Crée une diapositive NASA visuelle"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    title_shape = slide.shapes.title
    title_shape.text = "🚀 NASA Power of 10"
    
    rules = [
        ("1-2", "Simple\nControl Flow", RGBColor(52, 152, 219)),
        ("3", "No Dynamic\nMemory", RGBColor(231, 76, 60)),
        ("4", "Small\nFunctions", RGBColor(46, 204, 113)),
        ("5-6", "Check All\nReturns", RGBColor(241, 196, 15)),
        ("7-8", "Assertions +\nPointers", RGBColor(230, 126, 34)),
        ("9-10", "Zero Warnings\n+ Analysis", RGBColor(155, 89, 182)),
    ]
    
    positions = [
        (0.8, 2), (3.5, 2), (6.2, 2),
        (0.8, 4.5), (3.5, 4.5), (6.2, 4.5)
    ]
    
    for (rule, text, color), (left, top) in zip(rules, positions):
        add_diagram_box(slide, Inches(left), Inches(top), Inches(2.5), Inches(2), 
                       f"Règle {rule}\n\n{text}", color, 15)
    
    return slide

def create_comparison_slide(prs):
    """Crée une diapositive de comparaison visuelle"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    title_shape = slide.shapes.title
    title_shape.text = "❌ Avant vs ✅ Après"
    
    # Avant (gauche) - rouge
    add_diagram_box(slide, Inches(0.5), Inches(2), Inches(4.5), Inches(4.5),
                   "❌ AVANT\n\n" + 
                   "• Tout mélangé\n" +
                   "• Controller + DB\n" +
                   "• Impossible tester\n" +
                   "• Couplage fort\n" +
                   "• Code spaghetti",
                   RGBColor(192, 57, 43), 18)
    
    # Après (droite) - vert
    add_diagram_box(slide, Inches(5.5), Inches(2), Inches(4.5), Inches(4.5),
                   "✅ APRÈS\n\n" +
                   "• Séparation claire\n" +
                   "• Use Cases isolés\n" +
                   "• Testable facilement\n" +
                   "• Découplé\n" +
                   "• Clean Code",
                   RGBColor(39, 174, 96), 18)
    
    return slide

def create_ddd_building_blocks_slide(prs):
    """Crée une diapositive sur les building blocks DDD"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title_shape = slide.shapes.title
    title_shape.text = "DDD Building Blocks"
    
    blocks = [
        "🎯 Entity: Objet avec identité unique (User, Order)",
        "💎 Value Object: Défini par ses valeurs (Money, Email)",
        "📦 Aggregate: Cluster cohérent (Order + OrderLines)",
        "⚡ Domain Event: Événement métier (OrderPlaced)",
        "🔧 Domain Service: Logique multi-entités",
        "🏭 Factory: Création complexe d'objets",
        "📚 Repository: Abstraction de persistance"
    ]
    
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    
    for block in blocks:
        p = tf.add_paragraph()
        p.text = block
        p.font.size = Pt(18)
    
    return slide

def main():
    # Créer la présentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Titre
    create_title_slide(prs, 
                      "🚀 Workshop Architecture",
                      "Web & Embedded Systems")
    
    # Slide 2: Contenu
    create_content_slide(prs, "📋 Au Programme", [
        "🌐 Architecture Web",
        "",
        "🔧 Architecture C/Embedded",
        "",
        "💡 Bonnes Pratiques",
        "",
        "⚡ Patterns Avancés"
    ], 28)
    
    # === SECTION WEB ===
    
    # Slide 3: Clean Architecture
    create_title_slide(prs, "🌐 Clean Architecture")
    
    # Slide 4: Diagramme Architecture
    create_architecture_diagram_slide(prs)
    
    # Slide 5: SOLID
    create_solid_visual_slide(prs)
    
    # Slide 6: Avant/Après
    create_comparison_slide(prs)
    
    # Slide 7: DDD
    create_title_slide(prs, "📦 Domain-Driven Design")
    
    # Slide 8: DDD Visual
    create_ddd_visual_slide(prs)
    
    # Slide 9: Ubiquitous Language
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "🗣️ Ubiquitous Language"
    add_big_icon_text(slide, Inches(2.5), "💬", 
                     "Code = Langage Métier", 
                     RGBColor(52, 152, 219))
    
    # Slide 10: CQRS
    create_title_slide(prs, "⚡ CQRS + Event Sourcing")
    
    # Slide 11: CQRS Diagram
    create_cqrs_diagram_slide(prs)
    
    # Slide 12: Event Sourcing
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "📚 Event Sourcing"
    
    add_diagram_box(slide, Inches(1.5), Inches(2), Inches(7), Inches(1.2),
                   "Stocker l'HISTORIQUE pas l'ÉTAT",
                   RGBColor(52, 152, 219), 20)
    
    add_diagram_box(slide, Inches(2), Inches(3.5), Inches(6), Inches(3),
                   "AccountCreated\n⬇️\nMoneyDeposited +100€\n⬇️\nMoneyWithdrawn -50€\n⬇️\nBalance = 50€",
                   RGBColor(46, 204, 113), 18)
    
    # === SECTION C/EMBEDDED ===
    
    # Slide 13: C Section
    create_title_slide(prs, "🔧 Architecture C/Embedded")
    
    # Slide 14: NASA Rules Visual
    create_nasa_visual_slide(prs)
    
    # Slide 15: Why NASA
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "🚀 Pourquoi NASA Rules?"
    add_big_icon_text(slide, Inches(1.5), "🛡️", 
                     "Zéro Tolérance aux Bugs", 
                     RGBColor(231, 76, 60))
    add_big_icon_text(slide, Inches(4), "⚡", 
                     "Performance Prévisible", 
                     RGBColor(46, 204, 113))
    
    # Slide 16: Memory Safety
    create_title_slide(prs, "🛡️ Memory Safety")
    
    # Slide 17: Memory Hierarchy
    create_memory_hierarchy_slide(prs)
    
    # Slide 18: Static > All
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "⭐ Allocation Statique"
    add_big_icon_text(slide, Inches(2), "🏆", 
                     "STATIQUE > DYNAMIQUE", 
                     RGBColor(46, 204, 113))
    
    note = slide.shapes.add_textbox(Inches(2), Inches(5), Inches(6), Inches(1.5))
    note.text_frame.text = "int data[1000];\n\nPas de malloc() !\nPas de free() !\nZéro fuite !"
    note.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    note.text_frame.paragraphs[0].font.size = Pt(24)
    note.text_frame.paragraphs[0].font.bold = True
    
    # Slide 19: Layered C
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "🏗️ Architecture en Couches C"
    
    c_layers = [
        ("APPLICATION", RGBColor(52, 152, 219), 2),
        ("SERVICE", RGBColor(46, 204, 113), 3.2),
        ("DRIVER", RGBColor(241, 196, 15), 4.4),
        ("HAL (Hardware)", RGBColor(231, 76, 60), 5.6)
    ]
    
    for text, color, top in c_layers:
        add_diagram_box(slide, Inches(2.5), Inches(top), Inches(5), Inches(1), 
                       text, color, 20)
    
    # === PRATIQUE ===
    
    # Slide 20: Exercices
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "💪 Exercices Pratiques"
    
    add_diagram_box(slide, Inches(0.8), Inches(2.5), Inches(4), Inches(4),
                   "🌐 WEB\n\n" +
                   "6 exercices\nClean Arch\n\n" +
                   "DDD patterns\n\n" +
                   "CQRS projections",
                   RGBColor(52, 152, 219), 18)
    
    add_diagram_box(slide, Inches(5.2), Inches(2.5), Inches(4), Inches(4),
                   "🔧 C\n\n" +
                   "NASA Rules\n\n" +
                   "Arena allocator\n\n" +
                   "Drivers HAL",
                   RGBColor(231, 76, 60), 18)
    
    # Slide 21: Quand utiliser?
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "🤔 Quand Utiliser?"
    
    use_cases = [
        ("Clean Arch", "App complexe\n+6 mois", RGBColor(52, 152, 219), 1, 2.5),
        ("DDD", "Métier riche\nRègles business", RGBColor(46, 204, 113), 3.5, 2.5),
        ("CQRS", "Audit\nMulti-vues", RGBColor(241, 196, 15), 6, 2.5),
        ("NASA", "Critique\nZéro bug", RGBColor(231, 76, 60), 1, 5),
        ("Memory Safe", "Embedded\nTemps réel", RGBColor(230, 126, 34), 3.5, 5),
        ("⚠️", "Simple?\nPas de over-engineering!", RGBColor(149, 165, 166), 6, 5),
    ]
    
    for title, text, color, left, top in use_cases:
        add_diagram_box(slide, Inches(left), Inches(top), Inches(2.8), Inches(1.8),
                       f"{title}\n\n{text}", color, 14)
    
    # Slide 22: Anti-Patterns
    create_content_slide(prs, "⚠️ Anti-Patterns", [
        "❌ God Class",
        "❌ Anemic Domain",
        "❌ Big Ball of Mud",
        "❌ Magic Numbers",
        "❌ Ignorer les erreurs"
    ], 24)
    
    # Slide 23: Ressources
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "📚 Ressources"
    
    add_diagram_box(slide, Inches(1), Inches(2), Inches(8), Inches(1.2),
                   "📖 Clean Architecture - Robert C. Martin",
                   RGBColor(52, 152, 219), 18)
    
    add_diagram_box(slide, Inches(1), Inches(3.5), Inches(8), Inches(1.2),
                   "📖 Domain-Driven Design - Eric Evans",
                   RGBColor(46, 204, 113), 18)
    
    add_diagram_box(slide, Inches(1), Inches(5), Inches(8), Inches(1.2),
                   "🚀 NASA JPL Coding Standards",
                   RGBColor(231, 76, 60), 18)
    
    # Slide 24: Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "✨ Conclusion"
    add_big_icon_text(slide, Inches(2), "🎯", 
                     "Architecture = Fondation", 
                     RGBColor(52, 152, 219))
    add_big_icon_text(slide, Inches(4.5), "💪", 
                     "Practice Makes Perfect", 
                     RGBColor(46, 204, 113))
    
    # Slide 25: Questions
    create_title_slide(prs, "Questions? 💬", "Merci ! 🚀")
    
    # Sauvegarder
    output_file = "/home/hugo/Epitech/workshop-project-architecture/Workshop_Architecture_Presentation.pptx"
    prs.save(output_file)
    print(f"✅ Présentation créée: {output_file}")
    print(f"📊 Nombre de diapositives: {len(prs.slides)}")

if __name__ == "__main__":
    main()
